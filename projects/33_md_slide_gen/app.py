import io
import os
import re
import sys
from pathlib import Path

from dotenv import load_dotenv
from flask import Flask, render_template, request, send_file, flash, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).parents[2]))
from shared.csrf_setup import init_csrf
from shared.utils import call_claude_text

load_dotenv()

THEME_COLORS = {
    "blue": "1565C0",
    "green": "2E7D32",
    "orange": "EF6C00",
    "purple": "6A1B9A",
}

OPTIMIZE_PROMPT = (
    "次のMarkdownテキストを、プレゼンテーション用のスライド構成に整理してください。\n"
    "ルール:\n"
    "- スライドの区切りは `---` を1行で入れる\n"
    "- 各スライドの先頭行は `# タイトル` 形式の見出しにする\n"
    "- 本文は `- ` から始まる箇条書きにする（1スライドあたり3〜5項目程度）\n"
    "- Markdown本文のみを出力し、説明文やコードブロックの囲みは付けない\n\n"
    "元のテキスト:\n"
)


def _split_slides(markdown_text: str) -> list[str]:
    slides = re.split(r"^\s*---\s*$", markdown_text, flags=re.MULTILINE)
    return [s.strip() for s in slides if s.strip()]


def _parse_slide(slide_text: str) -> tuple[str, list[str]]:
    lines = [line.rstrip() for line in slide_text.splitlines() if line.strip()]
    title = ""
    bullets = []
    for line in lines:
        heading_match = re.match(r"^#+\s+(.*)", line)
        bullet_match = re.match(r"^[-*]\s+(.*)", line)
        if heading_match and not title:
            title = heading_match.group(1).strip()
        elif bullet_match:
            bullets.append(bullet_match.group(1).strip())
        elif not title:
            title = line.strip()
        else:
            bullets.append(line.strip())
    return title, bullets


def _build_pptx(markdown_text: str, theme: str) -> io.BytesIO:
    color_hex = THEME_COLORS.get(theme, THEME_COLORS["blue"])
    from pptx.dml.color import RGBColor

    prs = Presentation()
    title_layout = prs.slide_layouts[1]

    for slide_text in _split_slides(markdown_text):
        title, bullets = _parse_slide(slide_text)
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title or "（タイトル未設定）"
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(color_hex)

        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()
        for i, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(20)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def create_app():
    app = Flask(__name__)
    app.config["SECRET_KEY"] = os.environ.get("SECRET_KEY", os.urandom(24).hex())
    # hub_appの全アプリが同じCookie名"session"を共有すると、別アプリへの訪問で
    # セッションが上書きされCSRFトークンが消える問題があるため、一意な名前にする。
    app.config["SESSION_COOKIE_NAME"] = "session_33_md_slide_gen"
    init_csrf(app)

    @app.route("/", methods=["GET"])
    def index():
        return render_template("index.html", themes=THEME_COLORS.keys())

    @app.route("/preview", methods=["POST"])
    def preview():
        markdown_text = request.form.get("markdown", "").strip()
        theme = request.form.get("theme", "blue")
        optimize = request.form.get("optimize") == "on"

        if not markdown_text:
            flash("Markdownテキストを入力してください。", "error")
            return redirect(url_for("index"))

        if optimize and os.environ.get("GEMINI_API_KEY"):
            try:
                markdown_text = call_claude_text(
                    None, "claude-haiku-4-5-20251001", 2000, OPTIMIZE_PROMPT + markdown_text
                ).strip()
            except Exception:
                flash("AI構成最適化に失敗したため、元のテキストでプレビューします。", "error")

        return render_template("preview.html", markdown=markdown_text, theme=theme)

    @app.route("/export/pptx", methods=["POST"])
    def export_pptx():
        markdown_text = request.form.get("markdown", "").strip()
        theme = request.form.get("theme", "blue")

        if not markdown_text:
            flash("Markdownテキストを入力してください。", "error")
            return redirect(url_for("index"))

        pptx_buf = _build_pptx(markdown_text, theme)
        return send_file(
            pptx_buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name="slides.pptx",
        )

    @app.route("/portfolio")
    def portfolio():
        return redirect("https://ai-labo.space/", code=301)

    return app


app = create_app()

if __name__ == "__main__":
    app.run(debug=False, port=5033)
